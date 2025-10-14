from __future__ import annotations
from typing import Any, Dict, List
from .base import BaseAgent, AgentResponse
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt

ARCHIVE_DIR = os.path.join(os.path.dirname(__file__), "..", "reports", "archive")

class ReportGeneratorAgent(BaseAgent):
    name = "report"

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def _add_bullets_slide(self, prs: Presentation, title: str, bullets: List[str]) -> None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            if i == 0:
                body.text = b
            else:
                body.add_paragraph().text = b

    def _add_table_slide(self, prs: Presentation, title: str, headers: List[str], rows: List[List[str]]) -> None:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        cols = len(headers)
        rows_n = len(rows) + 1
        table = slide.shapes.add_table(rows_n, cols, Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.5)).table
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
        for r_i, r in enumerate(rows, start=1):
            for c_i, val in enumerate(r):
                table.cell(r_i, c_i).text = str(val)

    def run(self, story: Dict[str, Any], **kwargs: Any) -> AgentResponse:
        os.makedirs(ARCHIVE_DIR, exist_ok=True)
        prs = Presentation()

        molecule = story.get("molecule", "Unknown")
        indication = story.get("indication", "Unknown")
        timestamp = datetime.utcnow().strftime("%Y%m%d-%H%M%S")
        fname = f"innovation_story_{molecule}_{indication}_{timestamp}.pptx".replace(" ", "_")
        fpath = os.path.abspath(os.path.join(ARCHIVE_DIR, fname))

        # Slide 1: Title
        self._add_title_slide(
            prs,
            title=f"Innovation Opportunity: {molecule}",
            subtitle=f"Indication: {indication} | Generated {timestamp} UTC",
        )

        # Slide 2: Unmet needs
        unmet = story.get("unmet_needs", ["Synthesize unmet needs from internal and web sources."])
        self._add_bullets_slide(prs, "Unmet Medical Needs", unmet)

        # Slide 3: Clinical trials
        trials = story.get("clinical_trials", [])
        trials_rows: List[List[str]] = [
            [t.get("nct_id"), t.get("phase"), t.get("sponsor"), t.get("status")] for t in trials[:6]
        ]
        self._add_table_slide(prs, "Ongoing Clinical Trials (sample)", ["NCT ID", "Phase", "Sponsor", "Status"], trials_rows)

        # Slide 4: Repurposing hypotheses
        hypotheses = story.get("repurposing_hypotheses", [
            "Evaluate alternative dosage forms for pediatric adherence",
            "Explore comorbid indications with overlapping pathophysiology",
        ])
        self._add_bullets_slide(prs, "Repurposing/Alternate Use Hypotheses", hypotheses)

        # Slide 5: Patent landscape
        patents = story.get("patents", [])
        patent_rows: List[List[str]] = [
            [p.get("patent_id"), p.get("assignee"), p.get("expiry_year"), p.get("status"), p.get("fto_flag")] for p in patents[:6]
        ]
        self._add_table_slide(
            prs,
            "Patent Landscape (sample)",
            ["Patent ID", "Assignee", "Expiry", "Status", "FTO"],
            patent_rows,
        )

        prs.save(fpath)
        return AgentResponse({"file_path": fpath, "file_name": fname})
